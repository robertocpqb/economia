from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
 
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
 
# Helper function
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
 
# Slide 1
add_slide("Reforma Tributária", "Nível de preparação da indústria")